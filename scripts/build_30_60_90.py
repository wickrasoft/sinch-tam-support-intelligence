#!/usr/bin/env python3
"""Generate a Sinch-branded 30-60-90 Day Plan from the corporate template.

Framed for the EMEA TAM Team Lead (player-coach) role and mapped to the job
description: team leadership, performance management, escalation management,
stakeholder collaboration, operational excellence, reporting, and development.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

SRC = "/Users/kalwic/Documents/Resume/Kalum Wickramatunga/Kalum - Sinch - TAM/Sinch - Template.pptx"
OUT = "/Users/kalwic/dev/TAM Assignment/Kalum-Wickramatunga-30-60-90-Day-Plan.pptx"
ASSETS = os.path.join(os.path.dirname(__file__), "assets")
IMG_BALANCE = os.path.join(ASSETS, "balance.jpg")
IMG_TEAM = os.path.join(ASSETS, "team.jpg")
IMG_PITCREW = os.path.join(ASSETS, "pitcrew-onbrand.png")
IMG_COLLAB = os.path.join(ASSETS, "collab.jpg")

# Sinch brand palette (from template theme1.xml)
BLUE = RGBColor(0x18, 0x60, 0xF0)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
Y_LIGHT = RGBColor(0xFF, 0xF6, 0xB3)   # accent2
Y_MID = RGBColor(0xFF, 0xF1, 0x94)     # accent3
Y_DEEP = RGBColor(0xFF, 0xE9, 0x7A)    # accent4
GREY = RGBColor(0xF0, 0xF3, 0xF5)      # accent5
LINE = RGBColor(0xC8, 0xD2, 0xDC)
INK = RGBColor(0x1A, 0x1A, 0x1A)

HEAD = "Host Grotesk SemiBold"
BODY = "Host Grotesk"

prs = Presentation(SRC)
SW, SH = prs.slide_width, prs.slide_height

# ---- remove all existing slides (keep masters/layouts/theme/fonts) ----
sldIdLst = prs.slides._sldIdLst
for sldId in list(sldIdLst):
    rId = sldId.get(qn("r:id"))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


def layout(name):
    for m in prs.slide_masters:
        for l in m.slide_layouts:
            if l.name == name:
                return l
    raise KeyError(name)


def get_ph(slide, idx):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def remove_ph(slide, idx):
    ph = get_ph(slide, idx)
    if ph is not None:
        ph._element.getparent().remove(ph._element)


def strip_prompts(slide):
    for ph in list(slide.placeholders):
        t = ph.placeholder_format.type
        if t in (16, 15):  # 16 = DATE, 15 = FOOTER
            ph._element.getparent().remove(ph._element)


def set_text(ph, text, size=None, bold=None, color=None, font=BODY, align=None):
    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font
    if size is not None:
        f.size = Pt(size)
    if bold is not None:
        f.bold = bold
    if color is not None:
        f.color.rgb = color
    if align is not None:
        p.alignment = align
    return tf


def set_bullets(ph, items, size=16, color=INK):
    tf = ph.text_frame
    tf.word_wrap = True
    tf.clear()
    first = True
    for text, level in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        run = p.add_run()
        run.text = text
        run.font.name = BODY
        run.font.size = Pt(size - (1 if level > 0 else 0))
        run.font.color.rgb = color
        p.space_after = Pt(5)
    return tf


def add_box(slide, x, y, w, h, fill, line=None, radius=0.08):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    shp.shadow.inherit = False
    return shp


def para(tf, text, size, bold, color, font=BODY, space_before=0, space_after=4, first=False, align=PP_ALIGN.LEFT, level=0):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return p


def add_image_fit(slide, path, x, y, max_w, max_h):
    pic = slide.shapes.add_picture(path, x, y, width=max_w)
    if pic.height > max_h:
        ratio = int(max_h) / pic.height
        pic.height = int(max_h)
        pic.width = int(pic.width * ratio)
    pic.left = int(x + (int(max_w) - pic.width) / 2)
    pic.top = int(y + (int(max_h) - pic.height) / 2)
    return pic


def title_content(name_title, eyebrow):
    s = prs.slides.add_slide(layout("Title and content"))
    strip_prompts(s)
    set_text(get_ph(s, 19), name_title, size=26, bold=True, color=BLACK, font=HEAD)
    if eyebrow:
        set_text(get_ph(s, 13), eyebrow, size=14, color=BLUE, font=BODY)
    return s


def grouped_left(slide, x, y, w, h, groups):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for heading, lines in groups:
        para(tf, heading, 15, True, BLUE, font=HEAD, first=first, space_before=0 if first else 10, space_after=3)
        first = False
        for ln in lines:
            para(tf, ln, 13, False, INK, font=BODY, space_after=3, level=1)
    return tb


# =====================================================================
# SLIDE 1 — Cover
# =====================================================================
s = prs.slides.add_slide(layout("Title slide"))
strip_prompts(s)
remove_ph(s, 13)
set_text(get_ph(s, 0), "My First 90 Days", size=40, bold=True, color=BLACK, font=HEAD)
set_text(get_ph(s, 1), "A structured 30, 60 and 90 day plan", size=20, color=BLUE, font=HEAD)
set_text(get_ph(s, 14), "EMEA TEAM LEAD, TECHNICAL ACCOUNT MANAGEMENT", size=12, bold=True, color=BLUE, font=HEAD)
set_text(get_ph(s, 17), "Kalum Wickramatunga", size=14, bold=True, color=BLACK, font=HEAD)
set_text(get_ph(s, 19), "My plan to lead the EMEA team and manage my own account portfolio", size=11, color=INK, font=BODY)

# =====================================================================
# SLIDE 2 — Agenda
# =====================================================================
s = prs.slides.add_slide(layout("Agenda"))
strip_prompts(s)
set_text(get_ph(s, 0), "Agenda", size=28, bold=True, color=BLACK, font=HEAD)
set_bullets(get_ph(s, 17), [
    ("What I bring to the role", 0),
    ("My approach and the two sides of the role", 0),
    ("The first 90 days: 30, 60, and 90", 0),
    ("Leading and developing the team", 0),
    ("Managing critical escalations", 0),
    ("Improving how the team works", 0),
    ("Working with the wider Sinch team", 0),
    ("Measuring success and reporting", 0),
], size=18)

# =====================================================================
# SLIDE 3 — What I bring to the role (with image)
# =====================================================================
s = title_content("What I bring to the role", "Aligned to the requirements of the role")
remove_ph(s, 1)
grouped_left(s, Inches(0.6), Inches(2.0), Inches(7.0), Inches(4.7), [
    ("Experience", [
        "Several years leading and delivering in technical account management and enterprise support",
    ]),
    ("Industry and technical", [
        "A solid grounding in CPaaS and cloud communications, with hands on experience across SMS, 10DLC, Conversation API, Voice, and Numbers",
        "Comfortable steering a technical team through demanding enterprise problems",
    ]),
    ("Leadership and delivery", [
        "A history of building teams that perform consistently and keep customers loyal",
        "Experienced in owning high stakes escalations and working with senior stakeholders",
    ]),
    ("Ways of working", [
        "Analytical, steady under pressure, and comfortable when priorities shift quickly",
    ]),
])
add_image_fit(s, IMG_COLLAB, Inches(8.0), Inches(2.0), Inches(4.8), Inches(4.5))

# =====================================================================
# SLIDE 4 — Approach and principles
# =====================================================================
s = title_content("My approach and the principles I work by", "The principles that will guide my leadership from day one")
set_bullets(get_ph(s, 1), [
    ("Maintain a customer first focus, protecting adoption, retention, and trust across every account", 0),
    ("Listen before acting, ensuring I understand the accounts, the data, and the team before making changes", 0),
    ("Lead by example, remaining close to live accounts so that I coach from experience rather than theory", 0),
    ("Develop the people around me by providing clarity, support, and the space to perform at their best", 0),
    ("Make decisions based on data, using signals such as SLA, CSAT, response and resolution times, and reopened tickets", 0),
    ("Operate as one Sinch, partnering closely with Sales, Support, Solutions, Product, and the TOC", 0),
    ("Take ownership of outcomes, with clear commitments, consistent follow through, and no surprises", 0),
], size=16)

# =====================================================================
# SLIDE 5 — The two sides of the role (with image)
# =====================================================================
s = title_content("The two sides of the role", "Balancing individual contribution with team leadership")
remove_ph(s, 1)
tb = s.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(6.7), Inches(4.7))
tf = tb.text_frame
tf.word_wrap = True
para(tf, "Individual contribution", 16, True, BLUE, font=HEAD, first=True, space_after=4)
for ln in [
    "Manage my own portfolio of strategic and escalation accounts",
    "Remain hands on with QBRs, renewals, and technical issues",
    "Maintain product depth across the Sinch stack",
]:
    para(tf, ln, 14, False, INK, font=BODY, space_after=4, level=1)
para(tf, "Team leadership", 16, True, BLUE, font=HEAD, space_before=10, space_after=4)
for ln in [
    "Lead, coach, and support the EMEA TAM team",
    "Own team coverage, priorities, and ways of working",
    "Represent the team to leadership and remove blockers promptly",
]:
    para(tf, ln, 14, False, INK, font=BODY, space_after=4, level=1)
para(tf, "I anticipate dedicating approximately two thirds of my time to leading the team and one third to my own accounts, adjusting as the team requires.",
     13, False, INK, font=BODY, space_before=12, space_after=0)
add_image_fit(s, IMG_BALANCE, Inches(7.7), Inches(2.1), Inches(5.1), Inches(4.3))

# =====================================================================
# SLIDE 6 — At a glance (3 columns)
# =====================================================================
s = prs.slides.add_slide(layout("Only title"))
strip_prompts(s)
set_text(get_ph(s, 19), "How my first 90 days break down", size=26, bold=True, color=BLACK, font=HEAD)
cols = [
    ("30", "Learn and listen", Y_LIGHT, BLACK, [
        "Onboard and absorb",
        "Meet the team and review accounts",
        "Assess team and account health",
    ]),
    ("60", "Contribute and optimise", Y_DEEP, BLACK, [
        "Take ownership of accounts",
        "Coach and support the team",
        "Resolve key friction points",
    ]),
    ("90", "Lead and scale", BLUE, WHITE, [
        "Deliver QBRs and reviews",
        "Raise team performance",
        "Demonstrate measurable impact",
    ]),
]
left0, top0, gap, cw, ch = Inches(0.6), Inches(2.35), Inches(0.35), Inches(3.84), Inches(4.4)
for i, (num, title, fill, fg, lines) in enumerate(cols):
    x = Emu(int(left0) + i * (int(cw) + int(gap)))
    box = add_box(s, x, top0, cw, ch, fill)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.3)
    tf.margin_right = Inches(0.3)
    tf.margin_top = Inches(0.3)
    para(tf, num, 60, True, fg, font=HEAD, first=True, space_after=0)
    para(tf, "DAYS", 13, True, fg, font=HEAD, space_after=8)
    para(tf, title, 15, True, fg, font=HEAD, space_after=12)
    for ln in lines:
        para(tf, ln, 14, False, fg, font=BODY, space_after=6)
cap = s.shapes.add_textbox(left0, Emu(int(top0) + int(ch) + int(Inches(0.15))), Inches(12.1), Inches(0.4))
para(cap.text_frame,
     "Progressing from learning the team and the business, to taking ownership and coaching, to leading proactively with measurable results.",
     12, False, INK, font=BODY, first=True)

# =====================================================================
# Phase detail slides
# =====================================================================
def phase_slide(title, eyebrow, focus, actions, outcomes, accent, chip_label):
    s = title_content(title, eyebrow)
    chip = add_box(s, Inches(11.55), Inches(0.78), Inches(1.55), Inches(0.55), accent, radius=0.4)
    fgc = WHITE if accent == BLUE else BLACK
    para(chip.text_frame, chip_label, 15, True, fgc, font=HEAD, first=True, align=PP_ALIGN.CENTER)
    items = [("Focus", 0), (focus, 1), ("Key actions", 0)]
    for a in actions:
        items.append((a, 1))
    items.append(("Outcomes", 0))
    for o in outcomes:
        items.append((o, 1))
    tf = get_ph(s, 1).text_frame
    tf.word_wrap = True
    tf.clear()
    first = True
    headers = {"Focus", "Key actions", "Outcomes"}
    for text, level in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = level
        r = p.add_run()
        r.text = text
        r.font.name = HEAD if text in headers else BODY
        r.font.size = Pt(16 if text in headers else 13)
        r.font.bold = text in headers
        r.font.color.rgb = BLUE if text in headers else INK
        p.space_after = Pt(3 if level > 0 else 6)
        p.space_before = Pt(8 if text in headers and not first else 0)
    return s

# SLIDE 7 — Days 1 to 30
phase_slide(
    "Days 1 to 30: Learn and listen",
    "Onboarding, building trust, and understanding the team and accounts",
    "Develop a thorough understanding of the portfolio, products, people, and team, and establish a clear baseline for both account and team health.",
    [
        "Complete onboarding and security training, and obtain access to Zendesk, the CRM, and reporting dashboards",
        "Review my accounts in detail, including history, open tickets, SLAs, CSAT, and upcoming renewals",
        "Meet each member of the EMEA team individually to understand their accounts, strengths, and challenges",
        "Map how work currently flows across Support, Sales, Solutions, and the TOC, including escalation paths",
        "Strengthen my product knowledge across SMS, 10DLC, Conversation API, Voice, and Numbers",
        "Assess team workload, coverage, SLAs, and morale",
    ],
    [
        "A clear baseline of account and team health",
        "Established trust with the team and continuity of service for customers",
    ],
    Y_LIGHT, "Days 1 to 30",
)

# SLIDE 8 — Days 31 to 60
phase_slide(
    "Days 31 to 60: Contribute and optimise",
    "Taking ownership of accounts and leading the team in practice",
    "Take full ownership of my accounts and begin leading day to day, delivering the first tangible improvements.",
    [
        "Assume full ownership of my portfolio of strategic and escalation accounts",
        "Hold regular one to ones and begin coaching the team on live calls and escalations",
        "Agree the measures that matter most for the team, covering SLAs, CSAT, and escalations, and track them closely",
        "Identify key friction points, such as duplicate tickets and supplier escalation gaps, and begin addressing them",
        "Establish fair, clear priorities and coverage to ensure nothing is missed",
        "Partner with Sales on renewals and growth opportunities across the team's accounts",
    ],
    [
        "Actively managing my accounts and leading the team",
        "A supported team and more consistent service for customers",
    ],
    Y_DEEP, "Days 31 to 60",
)

# SLIDE 9 — Days 61 to 90
phase_slide(
    "Days 61 to 90: Lead and scale",
    "Leading proactively, raising team performance, and delivering measurable results",
    "Lead proactively, raise standards across the team, and deliver measurable results across the portfolio.",
    [
        "Deliver my own QBRs and support the team in delivering theirs to a consistent standard",
        "Establish a weekly cadence for reviewing account and team health",
        "Deliver at least one process or tooling improvement end to end, for example reopening Zendesk tickets when a supplier responds",
        "Agree development goals with each team member and establish a shared team playbook",
        "Build a simple dashboard and begin giving the Global Head of TAM a regular, honest view of where we stand",
        "Present results and a clear plan for the following quarter",
    ],
    [
        "A proactive, data led team delivering steady, measurable improvement",
        "Recognised as a leader who delivers through the team and across my own accounts",
    ],
    BLUE, "Days 61 to 90",
)

# =====================================================================
# SLIDE 10 — Leading and developing the team (with image)
# =====================================================================
s = title_content("Leading and developing the team", "Daily leadership, coaching, and development")
remove_ph(s, 1)
tb = s.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(6.7), Inches(4.7))
tf = tb.text_frame
tf.word_wrap = True
lead_points = [
    "Be present day to day and set the tone for a team that works together, takes ownership, and keeps improving",
    "Hold regular one to ones focused on development, blockers, and wellbeing",
    "Coach on live accounts and escalations so the learning carries into real work",
    "Make sure each person has what they need to do the job well, from system access to product training",
    "Agree clear growth goals with each person and help them take the next step in their career",
    "Give honest, timely feedback, recognise good work, and clear obstacles quickly",
]
for i, ln in enumerate(lead_points):
    para(tf, ln, 14, False, INK, font=BODY, first=(i == 0), space_after=8, level=1)
add_image_fit(s, IMG_PITCREW, Inches(7.7), Inches(2.5), Inches(5.2), Inches(3.4))

# =====================================================================
# SLIDE 11 — Managing critical escalations (with image)
# =====================================================================
s = title_content("Managing critical escalations", "The person people turn to first when an issue is serious")
remove_ph(s, 1)
tb = s.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(6.7), Inches(4.7))
tf = tb.text_frame
tf.word_wrap = True
esc_points = [
    "Be the single owner customers and colleagues turn to when a critical issue lands",
    "Quickly bring together the right people across TAM, Support, Solutions, Engineering, and the TOC",
    "Keep customers and senior stakeholders informed with clear, regular updates",
    "Drive a fast resolution, then run a blameless review to capture what we learned",
    "Feed recurring issues back into product, process, and the team playbook",
]
for i, ln in enumerate(esc_points):
    para(tf, ln, 14, False, INK, font=BODY, first=(i == 0), space_after=9, level=1)
add_image_fit(s, IMG_TEAM, Inches(7.8), Inches(2.3), Inches(5.0), Inches(3.6))
strip = ["Identify", "Mobilise", "Communicate", "Resolve", "Review"]
sx, sy, sw, sh2, sgap = Inches(0.6), Inches(6.45), Inches(2.3), Inches(0.55), Inches(0.18)
for i, label in enumerate(strip):
    x = Emu(int(sx) + i * (int(sw) + int(sgap)))
    b = add_box(s, x, sy, sw, sh2, BLUE if i == 0 else GREY, radius=0.4)
    para(b.text_frame, label, 12, True, WHITE if i == 0 else INK, font=HEAD, first=True, align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 12 — Operational excellence & continuous improvement (steps)
# =====================================================================
s = title_content("Improving how the team works", "Making day to day work simpler, faster, and easier to scale")
remove_ph(s, 1)
steps = [
    ("1", "Identify friction", Y_LIGHT, BLACK, "Use data and team input to identify what slows delivery, such as duplicate tickets or stale escalations"),
    ("2", "Redesign the process", Y_DEEP, BLACK, "Simplify the workflow with clear ownership and appropriate Zendesk triggers and automations"),
    ("3", "Implement safely", Y_MID, BLACK, "Pilot first, ensure no customer disruption, and bring the team along with the change"),
    ("4", "Measure and refine", BLUE, WHITE, "Track the impact on SLAs and CSAT, then refine and share effective practices"),
]
lx, ty, cgap, ccw, cch = Inches(0.6), Inches(2.2), Inches(0.3), Inches(2.95), Inches(3.7)
for i, (num, title, fill, fg, desc) in enumerate(steps):
    x = Emu(int(lx) + i * (int(ccw) + int(cgap)))
    box = add_box(s, x, ty, ccw, cch, fill)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.22)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.25)
    para(tf, num, 40, True, fg, font=HEAD, first=True, space_after=2)
    para(tf, title, 16, True, fg, font=HEAD, space_after=8)
    para(tf, desc, 12, False, fg, font=BODY, space_after=0)
cap = s.shapes.add_textbox(lx, Inches(6.15), Inches(12.1), Inches(0.7))
para(cap.text_frame,
     "I will also keep coverage and workload balanced, so the team spends its time where it makes the biggest difference for customers.",
     13, False, INK, font=BODY, first=True)

# =====================================================================
# SLIDE 13 — Cross-functional collaboration (hub and spoke)
# =====================================================================
s = title_content("Working with the wider Sinch team", "Keeping everyone aligned on what matters most to customers")
remove_ph(s, 1)


def connect(slide, p1, p2):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, p1[0], p1[1], p2[0], p2[1])
    c.line.color.rgb = LINE
    c.line.width = Pt(1.5)
    return c


def node(slide, cx, cy, w, h, text, fill, fg, size=12):
    x = Emu(int(cx) - int(w) // 2)
    y = Emu(int(cy) - int(h) // 2)
    box = add_box(slide, x, y, w, h, fill, radius=0.2)
    box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    para(box.text_frame, text, size, True, fg, font=HEAD, first=True, align=PP_ALIGN.CENTER)


center = (Inches(6.67), Inches(4.35))
sats = [
    (Inches(3.3), Inches(2.75), "Sales"),
    (Inches(10.0), Inches(2.75), "Service Operations"),
    (Inches(2.5), Inches(4.35), "Solutions"),
    (Inches(10.85), Inches(4.35), "Global Support"),
    (Inches(3.3), Inches(5.95), "Product & Engineering"),
    (Inches(10.0), Inches(5.95), "The TOC"),
]
for sx, sy, _ in sats:
    connect(s, center, (sx, sy))
for sx, sy, label in sats:
    node(s, sx, sy, Inches(2.5), Inches(0.8), label, GREY, INK, size=12)
node(s, center[0], center[1], Inches(2.9), Inches(1.0), "EMEA TAM team", BLUE, WHITE, size=15)
cap = s.shapes.add_textbox(Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.6))
para(cap.text_frame,
     "I will keep these teams pulling in the same direction for customers, and make sure the TAM team is well represented at leadership level.",
     13, False, INK, font=BODY, first=True, align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 14 — Measuring success and reporting
# =====================================================================
s = prs.slides.add_slide(layout("Only title"))
strip_prompts(s)
set_text(get_ph(s, 19), "Measuring success and reporting", size=26, bold=True, color=BLACK, font=HEAD)
metric_cards = [
    ("Customer", Y_LIGHT, BLACK, [
        "CSAT trend, and NPS where available",
        "Renewals and retention",
        "Customer references and advocacy",
    ]),
    ("Service quality", Y_MID, BLACK, [
        "SLA compliance",
        "Improved response and resolution times",
        "Reduction in P1 and P2 escalations",
        "Reduction in stale and ageing tickets",
    ]),
    ("Team", Y_DEEP, BLACK, [
        "Team CSAT and SLA performance",
        "Balanced workload and coverage",
        "Team development and retention",
    ]),
    ("Growth", BLUE, WHITE, [
        "Growth opportunities identified with Sales",
        "Improvements delivered end to end",
        "Completion of my onboarding and certifications",
    ]),
]
lx, ty, cgap, ccw, cch = Inches(0.6), Inches(2.2), Inches(0.3), Inches(2.95), Inches(3.85)
for i, (title, fill, fg, lines) in enumerate(metric_cards):
    x = Emu(int(lx) + i * (int(ccw) + int(cgap)))
    box = add_box(s, x, ty, ccw, cch, fill)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.22)
    para(tf, title, 17, True, fg, font=HEAD, first=True, space_after=9)
    for ln in lines:
        para(tf, ln, 12, False, fg, font=BODY, space_after=6)
rep = add_box(s, lx, Inches(6.35), Inches(12.13), Inches(0.62), GREY, radius=0.18)
rep.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
para(rep.text_frame,
     "Each month I will give the Global Head of TAM a clear picture of how the team and our customers are doing, what the trends show, and the actions I am taking in response.",
     12, False, INK, font=BODY, first=True, align=PP_ALIGN.CENTER)

# =====================================================================
# SLIDE 15 — Operating cadence
# =====================================================================
s = title_content("Operating cadence", "The cadence I will maintain across the team, Sinch, and my accounts")
cph = get_ph(s, 1)
cx, cy, cwd, chd = cph.left, cph.top, cph.width, cph.height
cph._element.getparent().remove(cph._element)
data = [
    ("Frequency", "Focus"),
    ("Daily", "Monitor my portfolio and remain available to the team for live issues"),
    ("Weekly", "Team one to ones and a team sync, alongside account health reviews and updates with Sales and Support"),
    ("Bi-weekly", "Review escalation and ticket trends, and align with Service Operations on priorities"),
    ("Monthly", "A one to one with my manager, team performance and trend reporting, and account success plan updates"),
    ("Quarterly", "QBRs and service reviews with customers, with team development and planning for the next quarter"),
    ("Yearly", "Annual performance and development reviews, portfolio strategy, and goal setting for the year ahead"),
]
rows, colsN = len(data), 2
tbl_shape = s.shapes.add_table(rows, colsN, cx, cy, cwd, Inches(4.9))
table = tbl_shape.table
table.columns[0].width = Inches(2.5)
table.columns[1].width = Emu(int(cwd) - int(Inches(2.5)))
for r in range(rows):
    for c in range(colsN):
        cell = table.cell(r, c)
        cell.margin_left = Inches(0.12)
        cell.margin_right = Inches(0.12)
        cell.margin_top = Inches(0.04)
        cell.margin_bottom = Inches(0.04)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = data[r][c]
        run.font.name = HEAD if r == 0 else BODY
        run.font.size = Pt(14 if r == 0 else 12)
        run.font.bold = (r == 0) or (c == 0)
        if r == 0:
            run.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE
        else:
            run.font.color.rgb = INK
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else GREY

# =====================================================================
# SLIDE 16 — End cover
# =====================================================================
s = prs.slides.add_slide(layout("End cover"))
strip_prompts(s)
remove_ph(s, 13)
# End cover inherits a dark master background, so text is white with a Sinch
# yellow accent on the subtitle (the standard Sinch styling for dark slides).
set_text(get_ph(s, 0), "Thank you", size=34, bold=True, color=WHITE, font=HEAD)
set_text(get_ph(s, 1), "I look forward to leading and delivering in EMEA", size=18, bold=True, color=Y_DEEP, font=HEAD)
body = get_ph(s, 14)
if body is not None:
    set_text(body, "Kalum Wickramatunga", size=12, bold=True, color=WHITE, font=BODY)

prs.save(OUT)
print("Saved:", OUT)
print("Slides:", len(prs.slides._sldIdLst))
